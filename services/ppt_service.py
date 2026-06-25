"""
Scholarly — Professional PowerPoint report generator.
Uses python-pptx to create a 16:9 deck with branded slides.
"""
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── Brand palette ────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1b, 0x2e, 0x4b)
BLUE    = RGBColor(0x4a, 0x7d, 0xff)
GREEN   = RGBColor(0x58, 0xcc, 0x02)
RED     = RGBColor(0xff, 0x4b, 0x4b)
AMBER   = RGBColor(0xf5, 0xa6, 0x23)
WHITE   = RGBColor(0xff, 0xff, 0xff)
OFFWHITE= RGBColor(0xf4, 0xf4, 0xf4)
LGRAY   = RGBColor(0xe5, 0xe5, 0xe5)
MGRAY   = RGBColor(0xaf, 0xaf, 0xaf)
DGRAY   = RGBColor(0x3c, 0x3c, 0x3c)

W = Inches(13.33)
H = Inches(7.5)


def _rgb(r, g, b):
    return RGBColor(r, g, b)


# ── Slide helpers ────────────────────────────────────────────────────────────

def _fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    _fill_solid(shape, color)
    shape.line.fill.background()
    return shape


def _txb(slide, text, x, y, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def _stat_box(slide, x, y, w, h, value, label, val_color=NAVY, bg=OFFWHITE):
    box = _add_rect(slide, x, y, w, h, bg)
    box.line.color.rgb = LGRAY
    box.line.width = Pt(1)
    _txb(slide, str(value), x, y + Inches(0.35), w, Inches(0.7),
         28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    _txb(slide, label, x, y + Inches(1.0), w, Inches(0.45),
         10, bold=False, color=MGRAY, align=PP_ALIGN.CENTER)


def _section_label(slide, text, x, y, w=Inches(11)):
    _txb(slide, text.upper(), x, y, w, Inches(0.3),
         8, bold=True, color=MGRAY, align=PP_ALIGN.LEFT)


def _divider(slide, x, y, w):
    ln = _add_rect(slide, x, y, w, Pt(1.5), LGRAY)
    return ln


def _badge(slide, text, x, y, color):
    box = _add_rect(slide, x, y, Inches(0.85), Inches(0.28), color)
    box.line.fill.background()
    _txb(slide, text, x, y, Inches(0.85), Inches(0.28),
         8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Slide 1 — Cover ──────────────────────────────────────────────────────────

def _slide_cover(prs, title, subtitle, batch_name, section, date_str):
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)

    # Full dark background
    bg = _add_rect(slide, 0, 0, W, H, NAVY)

    # Accent bar left
    _add_rect(slide, 0, 0, Inches(0.12), H, BLUE)

    # Decorative circles (top-right)
    for i, (cx, cy, cr, alpha) in enumerate([
        (Inches(12.4), Inches(-0.4), Inches(2.2), 0.08),
        (Inches(11.5), Inches(0.8),  Inches(1.4), 0.05),
    ]):
        circ = slide.shapes.add_shape(9, cx - cr, cy - cr, cr*2, cr*2)  # oval
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE
        circ.line.fill.background()

    # Scholarly wordmark
    _txb(slide, "SCHOLARLY", Inches(0.5), Inches(0.5), Inches(5), Inches(0.5),
         11, bold=True, color=BLUE, align=PP_ALIGN.LEFT)

    # Thin rule under logo
    _add_rect(slide, Inches(0.5), Inches(1.05), Inches(12.3), Pt(1), _rgb(74,125,255))

    # Main title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r  = p.add_run(); r.text = title
    r.font.size = Pt(44); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Calibri"

    # Subtitle
    _txb(slide, subtitle, Inches(0.5), Inches(3.35), Inches(9), Inches(0.55),
         18, bold=False, color=_rgb(180,196,230), align=PP_ALIGN.LEFT)

    # Meta pills row
    meta_items = []
    if batch_name: meta_items.append(f"  {batch_name}  ")
    if section:    meta_items.append(f"  Section {section}  ")
    meta_items.append(f"  {date_str}  ")

    pill_x = Inches(0.5)
    for item in meta_items:
        w_approx = Inches(max(1.4, len(item) * 0.09))
        box = _add_rect(slide, pill_x, Inches(4.1), w_approx, Inches(0.38), _rgb(35,55,90))
        box.line.color.rgb = _rgb(74,125,255); box.line.width = Pt(0.75)
        _txb(slide, item.strip(), pill_x + Inches(0.1), Inches(4.12),
             w_approx - Inches(0.1), Inches(0.35),
             10, bold=False, color=_rgb(180,200,240), align=PP_ALIGN.LEFT)
        pill_x += w_approx + Inches(0.12)

    # Bottom rule
    _add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), BLUE)

    # Footer
    _txb(slide, "Confidential — Academic Performance Report",
         Inches(0.5), H - Inches(0.45), Inches(10), Inches(0.35),
         8, bold=False, color=MGRAY, align=PP_ALIGN.LEFT)

    return slide


# ── Slide 2 — Executive Summary ──────────────────────────────────────────────

def _slide_summary(prs, student_data, quiz_data, batch_name, section):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, WHITE)
    _add_rect(slide, 0, 0, Inches(0.06), H, BLUE)
    _add_rect(slide, 0, 0, W, Inches(1.1), NAVY)

    _txb(slide, "Executive Summary", Inches(0.25), Inches(0.22),
         Inches(10), Inches(0.6), 22, bold=True, color=WHITE)

    label = batch_name or "All Batches"
    if section: label += f" · Section {section}"
    _txb(slide, label, Inches(0.25), Inches(0.72), Inches(10), Inches(0.35),
         10, bold=False, color=_rgb(160,180,220))

    n = len(student_data)
    att_pcts = [s["attendance_pct"] for s in student_data]
    avg_att  = round(sum(att_pcts) / n) if n else 0
    at_risk  = sum(1 for p in att_pcts if p < 75)
    quiz_avgs= [s["quiz_avg"] for s in student_data if s["quiz_avg"] is not None]
    avg_quiz = round(sum(quiz_avgs) / len(quiz_avgs)) if quiz_avgs else None

    BOX_W = Inches(2.8)
    BOX_H = Inches(1.6)
    GAP   = Inches(0.26)
    START = Inches(0.55)
    ROW_Y = Inches(1.5)

    stats = [
        (str(n),               "Total Students",     NAVY),
        (f"{avg_att}%",        "Avg Attendance",     GREEN if avg_att >= 75 else RED),
        (f"{avg_quiz}%" if avg_quiz is not None else "—",
                               "Avg Quiz Score",     BLUE),
        (str(at_risk),         "Below 75% Attendance", RED if at_risk > 0 else GREEN),
    ]
    for i, (val, lbl, col) in enumerate(stats):
        bx = START + i * (BOX_W + GAP)
        box = _add_rect(slide, bx, ROW_Y, BOX_W, BOX_H, OFFWHITE)
        box.line.color.rgb = LGRAY; box.line.width = Pt(1)
        # colored top accent
        _add_rect(slide, bx, ROW_Y, BOX_W, Inches(0.07), col)
        _txb(slide, val, bx, ROW_Y + Inches(0.25), BOX_W, Inches(0.75),
             32, bold=True, color=col, align=PP_ALIGN.CENTER)
        _txb(slide, lbl, bx, ROW_Y + Inches(0.95), BOX_W, Inches(0.45),
             10, color=MGRAY, align=PP_ALIGN.CENTER)

    # Key insights section
    _section_label(slide, "Key Insights", Inches(0.55), Inches(3.35))
    _divider(slide, Inches(0.55), Inches(3.62), Inches(12.2))

    insights = []
    if avg_att >= 90:
        insights.append(("✓", f"Excellent class attendance at {avg_att}%", GREEN))
    elif avg_att >= 75:
        insights.append(("●", f"Attendance is acceptable at {avg_att}% — room to improve", AMBER))
    else:
        insights.append(("✗", f"Class attendance is critical at {avg_att}% — immediate action needed", RED))

    if at_risk > 0:
        insights.append(("⚠", f"{at_risk} student{'s' if at_risk>1 else ''} below 75% attendance threshold", RED))

    if avg_quiz is not None:
        if avg_quiz >= 80:
            insights.append(("✓", f"Strong quiz performance — class average {avg_quiz}%", GREEN))
        elif avg_quiz >= 60:
            insights.append(("●", f"Average quiz performance at {avg_quiz}% — targeted revision recommended", AMBER))
        else:
            insights.append(("✗", f"Low quiz performance at {avg_quiz}% — curriculum review recommended", RED))

    if len(quiz_data) > 0:
        insights.append(("📊", f"{len(quiz_data)} quiz{'zes' if len(quiz_data)>1 else ''} conducted this term", BLUE))

    for i, (icon, text, col) in enumerate(insights[:4]):
        y = Inches(3.75) + i * Inches(0.72)
        _add_rect(slide, Inches(0.55), y + Inches(0.06), Inches(0.28), Inches(0.28), col)
        _txb(slide, icon, Inches(0.55), y + Inches(0.03), Inches(0.28), Inches(0.32),
             9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, text, Inches(0.95), y, Inches(11.5), Inches(0.4),
             11, color=DGRAY)

    _add_rect(slide, 0, H - Inches(0.06), W, Inches(0.06), BLUE)
    return slide


# ── Slide 3 — Attendance Chart ───────────────────────────────────────────────

def _slide_attendance(prs, student_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, WHITE)
    _add_rect(slide, 0, 0, Inches(0.06), H, BLUE)
    _add_rect(slide, 0, 0, W, Inches(1.1), NAVY)

    _txb(slide, "Attendance Overview", Inches(0.25), Inches(0.22),
         Inches(10), Inches(0.6), 22, bold=True, color=WHITE)
    _txb(slide, "Student-wise attendance percentage",
         Inches(0.25), Inches(0.72), Inches(10), Inches(0.35),
         10, color=_rgb(160,180,220))

    if not student_data:
        _txb(slide, "No data available", Inches(3), Inches(3.5), Inches(7), Inches(1),
             18, color=MGRAY, align=PP_ALIGN.CENTER)
        return slide

    # Bar chart (top 15 students for readability)
    data = sorted(student_data, key=lambda s: s["attendance_pct"])[:20]
    cd   = ChartData()
    cd.categories = [s["name"].split()[0] for s in data]
    cd.add_series("Attendance %", [s["attendance_pct"] for s in data])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.45), Inches(1.25),
        Inches(8.5), Inches(5.8),
        cd
    ).chart

    chart.has_legend = False
    chart.has_title  = False
    plot  = chart.plots[0]
    ser   = plot.series[0]

    # Colour each bar green/amber/red
    from pptx.oxml.ns import qn
    from lxml import etree
    for i, s in enumerate(data):
        pct = s["attendance_pct"]
        col = GREEN if pct >= 90 else (AMBER if pct >= 75 else RED)
        pt  = ser.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = col

    # 75% threshold label
    _txb(slide, "75% Threshold", Inches(9.1), Inches(1.3), Inches(3.5), Inches(0.4),
         9, bold=True, color=RED)

    # Legend
    LY = Inches(1.4)
    for col, label in [(GREEN, "≥ 90% Excellent"), (AMBER, "75–89% Acceptable"), (RED, "< 75% At Risk")]:
        _add_rect(slide, Inches(9.1), LY + Inches(0.04), Inches(0.22), Inches(0.22), col)
        _txb(slide, label, Inches(9.4), LY, Inches(3.5), Inches(0.32), 9, color=DGRAY)
        LY += Inches(0.38)

    # Stats sidebar
    att_pcts = [s["attendance_pct"] for s in student_data]
    n = len(att_pcts)
    perfect  = sum(1 for p in att_pcts if p == 100)
    low      = sum(1 for p in att_pcts if p < 75)
    avg      = round(sum(att_pcts)/n) if n else 0

    LY = Inches(2.5)
    for val, lbl, col in [
        (f"{avg}%",   "Class Avg",  NAVY),
        (str(perfect),"100% Att.",  GREEN),
        (str(low),    "At Risk",    RED),
    ]:
        _stat_box(slide, Inches(9.1), LY, Inches(3.8), Inches(1.3), val, lbl, col, OFFWHITE)
        LY += Inches(1.45)

    _add_rect(slide, 0, H - Inches(0.06), W, Inches(0.06), BLUE)
    return slide


# ── Slide 4 — Quiz Performance Chart ─────────────────────────────────────────

def _slide_quizzes(prs, quiz_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, WHITE)
    _add_rect(slide, 0, 0, Inches(0.06), H, BLUE)
    _add_rect(slide, 0, 0, W, Inches(1.1), NAVY)

    _txb(slide, "Quiz Performance", Inches(0.25), Inches(0.22),
         Inches(10), Inches(0.6), 22, bold=True, color=WHITE)
    _txb(slide, "Average score per quiz conducted",
         Inches(0.25), Inches(0.72), Inches(10), Inches(0.35),
         10, color=_rgb(160,180,220))

    if not quiz_data:
        _txb(slide, "No quizzes conducted yet", Inches(3), Inches(3.5), Inches(7), Inches(1),
             18, color=MGRAY, align=PP_ALIGN.CENTER)
        return slide

    quizzes_with_data = [q for q in quiz_data if q.get("avg_score") is not None]

    if quizzes_with_data:
        cd = ChartData()
        cd.categories = [q["title"][:20] for q in quizzes_with_data]
        cd.add_series("Avg Score %", [q["avg_score"] for q in quizzes_with_data])

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.45), Inches(1.25),
            Inches(8.5), Inches(5.0),
            cd
        ).chart

        chart.has_legend = False
        chart.has_title  = False

        ser = chart.plots[0].series[0]
        for i, q in enumerate(quizzes_with_data):
            avg = q["avg_score"]
            col = GREEN if avg >= 80 else (AMBER if avg >= 60 else RED)
            pt  = ser.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = col

    # Quiz table below chart
    TABLE_Y = Inches(6.45)
    _txb(slide, "QUIZ SUMMARY", Inches(0.45), TABLE_Y - Inches(0.28), Inches(8), Inches(0.25),
         7, bold=True, color=MGRAY)
    _divider(slide, Inches(0.45), TABLE_Y - Inches(0.05), Inches(8.5))

    # Stats sidebar
    avgs = [q["avg_score"] for q in quiz_data if q.get("avg_score") is not None]
    SY = Inches(1.35)
    for val, lbl, col in [
        (str(len(quiz_data)),                        "Total Quizzes",  NAVY),
        (f"{round(sum(avgs)/len(avgs))}%" if avgs else "—", "Overall Avg", BLUE),
        (f"{max(avgs)}%" if avgs else "—",           "Best Quiz Avg",  GREEN),
        (f"{min(avgs)}%" if avgs else "—",           "Lowest Quiz Avg",RED),
    ]:
        _stat_box(slide, Inches(9.1), SY, Inches(3.8), Inches(1.28), val, lbl, col, OFFWHITE)
        SY += Inches(1.42)

    _add_rect(slide, 0, H - Inches(0.06), W, Inches(0.06), BLUE)
    return slide


# ── Slide 5 — Student Table ───────────────────────────────────────────────────

def _slide_student_table(prs, student_data, page=1, total_pages=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, WHITE)
    _add_rect(slide, 0, 0, Inches(0.06), H, BLUE)
    _add_rect(slide, 0, 0, W, Inches(1.1), NAVY)

    heading = "Student Performance Table"
    if total_pages > 1:
        heading += f"  ({page}/{total_pages})"
    _txb(slide, heading, Inches(0.25), Inches(0.22),
         Inches(10), Inches(0.6), 22, bold=True, color=WHITE)
    _txb(slide, "Attendance percentage, quiz average and performance trend",
         Inches(0.25), Inches(0.72), Inches(10), Inches(0.35),
         10, color=_rgb(160,180,220))

    COLS  = ["#", "Name", "Batch / Section", "Attendance", "Quiz Avg", "Quizzes", "Trend"]
    WIDTHS= [Inches(0.4), Inches(3.0), Inches(2.2), Inches(1.5), Inches(1.5), Inches(1.2), Inches(1.5)]
    COL_X = [Inches(0.35)]
    for w in WIDTHS[:-1]:
        COL_X.append(COL_X[-1] + w)

    HDR_Y = Inches(1.22)
    HDR_H = Inches(0.38)
    _add_rect(slide, Inches(0.3), HDR_Y, Inches(12.7), HDR_H, NAVY)
    for j, (col, x) in enumerate(zip(COLS, COL_X)):
        _txb(slide, col, x, HDR_Y + Inches(0.06), WIDTHS[j], HDR_H - Inches(0.06),
             8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    ROW_H = Inches(0.46)
    for i, s in enumerate(student_data):
        ry  = HDR_Y + HDR_H + i * ROW_H
        bg  = OFFWHITE if i % 2 == 0 else WHITE
        _add_rect(slide, Inches(0.3), ry, Inches(12.7), ROW_H, bg)

        att  = s["attendance_pct"]
        quiz = s.get("quiz_avg")
        trend= s.get("trend")

        att_col  = GREEN if att  >= 90 else (AMBER if att  >= 75 else RED)
        quiz_col = (GREEN if quiz >= 80 else (AMBER if quiz >= 60 else RED)) if quiz is not None else MGRAY
        t_sym, t_col = ("▲", GREEN) if (trend or 0) > 2 else (("▼", RED) if (trend or 0) < -2 else ("—", MGRAY))

        row_vals = [
            str(i + 1),
            s["name"],
            f"{s.get('batch_name') or '—'} / {s.get('section') or '—'}",
            f"{att}%",
            f"{quiz}%" if quiz is not None else "—",
            str(s.get("quizzes_taken", 0)),
            t_sym,
        ]
        row_colors = [DGRAY, DGRAY, MGRAY, att_col, quiz_col, DGRAY, t_col]

        for j, (val, col, x) in enumerate(zip(row_vals, row_colors, COL_X)):
            _txb(slide, val, x, ry + Inches(0.09), WIDTHS[j], ROW_H - Inches(0.09),
                 9, bold=(j == 1), color=col,
                 align=PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT)

    _add_rect(slide, 0, H - Inches(0.06), W, Inches(0.06), BLUE)
    return slide


# ── Slide — Thank You / End ───────────────────────────────────────────────────

def _slide_end(prs, batch_name, section):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, W, H, NAVY)
    _add_rect(slide, 0, 0, Inches(0.12), H, BLUE)

    # Big decorative circle
    for cx, cy, cr in [(Inches(10.5), Inches(5.5), Inches(3.5)),
                       (Inches(11.8), Inches(2.8), Inches(1.8))]:
        c = slide.shapes.add_shape(9, cx - cr, cy - cr, cr*2, cr*2)
        c.fill.solid(); c.fill.fore_color.rgb = _rgb(30,50,80)
        c.line.fill.background()

    _txb(slide, "SCHOLARLY", Inches(0.5), Inches(0.5), Inches(5), Inches(0.5),
         11, bold=True, color=BLUE)
    _add_rect(slide, Inches(0.5), Inches(1.05), Inches(12.3), Pt(1), BLUE)

    _txb(slide, "Thank You", Inches(0.5), Inches(1.9), Inches(11), Inches(1.8),
         52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _txb(slide, "This report was auto-generated by Scholarly ERP.",
         Inches(0.5), Inches(3.65), Inches(9), Inches(0.55),
         14, color=_rgb(160,180,220))

    label_parts = []
    if batch_name: label_parts.append(batch_name)
    if section:    label_parts.append(f"Section {section}")
    if label_parts:
        _txb(slide, " · ".join(label_parts),
             Inches(0.5), Inches(4.25), Inches(9), Inches(0.5),
             11, color=BLUE)

    _txb(slide, datetime.now().strftime("Generated on %d %B %Y"),
         Inches(0.5), Inches(4.75), Inches(9), Inches(0.4),
         9, color=MGRAY)

    _add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), BLUE)
    return slide


# ── Main entry point ──────────────────────────────────────────────────────────

def generate_report_pptx(
    student_data: list[dict],
    quiz_data: list[dict],
    batch_name: str | None = None,
    section: str | None = None,
) -> bytes:
    """
    student_data: [{name, batch_name, section, attendance_pct, quiz_avg, quizzes_taken, trend}]
    quiz_data:    [{title, subject, created_at, avg_score, highest, lowest, completion_rate}]
    Returns raw bytes of the .pptx file.
    """
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    date_str = datetime.now().strftime("%B %Y")
    label    = batch_name or "All Classes"
    subtitle = f"Academic Performance Report — {label}"
    if section:
        subtitle += f" · Section {section}"

    _slide_cover(prs, "Class Performance Report", subtitle, batch_name, section, date_str)
    _slide_summary(prs, student_data, quiz_data, batch_name, section)
    _slide_attendance(prs, student_data)
    _slide_quizzes(prs, quiz_data)

    # Paginate student table (14 rows per slide)
    PAGE_SIZE = 14
    pages     = [student_data[i:i+PAGE_SIZE] for i in range(0, max(len(student_data), 1), PAGE_SIZE)]
    for pg, chunk in enumerate(pages, 1):
        _slide_student_table(prs, chunk, pg, len(pages))

    _slide_end(prs, batch_name, section)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
